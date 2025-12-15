"""
Table-Aware Document Extraction

Provides enhanced extraction for PDF and DOCX files with:
- Table detection and preservation
- Markdown table formatting
- MarkItDown integration (Microsoft's converter)
- Fallback strategies for compatibility
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from io import BytesIO

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Result of document extraction."""
    text: str
    tables: List[Dict[str, Any]] = field(default_factory=list)
    page_count: int = 1
    extraction_method: str = "unknown"


# ============================================================================
# PDF Extraction
# ============================================================================

def extract_pdf_with_tables(
    file_path: Union[str, Path, BytesIO],
    prefer_markitdown: bool = True
) -> ExtractionResult:
    """
    Extract text from PDF with table awareness.

    Uses a multi-strategy approach:
    1. MarkItDown (if available) - best for structure preservation
    2. pdfplumber with table extraction
    3. Basic pypdf fallback

    Args:
        file_path: Path to PDF file or BytesIO object
        prefer_markitdown: Try MarkItDown first (recommended)

    Returns:
        ExtractionResult with text, tables, and metadata
    """
    # Try MarkItDown first
    if prefer_markitdown:
        result = _try_markitdown(file_path)
        if result:
            return result

    # Try pdfplumber with tables
    result = _try_pdfplumber(file_path)
    if result:
        return result

    # Fallback to basic pypdf
    return _try_pypdf(file_path)


def _try_markitdown(file_path: Union[str, Path, BytesIO]) -> Optional[ExtractionResult]:
    """Try extraction with MarkItDown."""
    try:
        from markitdown import MarkItDown

        md = MarkItDown()

        # Handle BytesIO vs file path
        if isinstance(file_path, BytesIO):
            # MarkItDown needs a file path, so save temporarily
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                tmp.write(file_path.read())
                tmp_path = tmp.name
            file_path.seek(0)  # Reset for other extractors
            result = md.convert(tmp_path)
            Path(tmp_path).unlink()  # Clean up
        else:
            result = md.convert(str(file_path))

        if result and result.text_content:
            # Extract tables from markdown (tables are in the text as markdown)
            tables = _extract_markdown_tables(result.text_content)

            logger.info(f"MarkItDown extracted {len(tables)} tables")
            return ExtractionResult(
                text=result.text_content,
                tables=tables,
                page_count=1,  # MarkItDown doesn't report page count
                extraction_method="markitdown"
            )

    except ImportError:
        logger.debug("MarkItDown not available")
    except Exception as e:
        logger.warning(f"MarkItDown extraction failed: {e}")

    return None


def _try_pdfplumber(file_path: Union[str, Path, BytesIO]) -> Optional[ExtractionResult]:
    """Try extraction with pdfplumber."""
    try:
        import pdfplumber

        text_parts = []
        tables = []
        page_count = 0

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages, 1):
                # Extract tables first
                page_tables = page.extract_tables()
                for i, table in enumerate(page_tables):
                    if table and len(table) > 1:
                        markdown = _table_to_markdown(table)
                        tables.append({
                            "markdown": markdown,
                            "page": page_num,
                            "index": i
                        })

                # Extract text
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(page_text)

        # Combine text with tables
        full_text = "\n\n".join(text_parts)
        for table in tables:
            full_text += f"\n\n{table['markdown']}\n"

        logger.info(f"pdfplumber extracted {len(tables)} tables from {page_count} pages")
        return ExtractionResult(
            text=full_text,
            tables=tables,
            page_count=page_count,
            extraction_method="pdfplumber"
        )

    except ImportError:
        logger.debug("pdfplumber not available")
    except Exception as e:
        logger.warning(f"pdfplumber extraction failed: {e}")

    return None


def _try_pypdf(file_path: Union[str, Path, BytesIO]) -> ExtractionResult:
    """Basic extraction with pypdf (fallback)."""
    try:
        import pypdf

        text_parts = []
        page_count = 0

        # Handle BytesIO vs file path
        if isinstance(file_path, BytesIO):
            reader = pypdf.PdfReader(file_path)
        else:
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)

        page_count = len(reader.pages)

        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)

        logger.info(f"pypdf extracted {page_count} pages (basic mode)")
        return ExtractionResult(
            text="\n\n".join(text_parts),
            tables=[],
            page_count=page_count,
            extraction_method="pypdf_basic"
        )

    except ImportError:
        logger.error("pypdf not installed")
    except Exception as e:
        logger.error(f"pypdf extraction failed: {e}")

    return ExtractionResult(
        text="[PDF extraction failed]",
        tables=[],
        page_count=0,
        extraction_method="failed"
    )


# ============================================================================
# DOCX Extraction
# ============================================================================

def extract_docx_with_tables(
    file_path: Union[str, Path, BytesIO],
    prefer_markitdown: bool = True
) -> ExtractionResult:
    """
    Extract text from Word document with table awareness.

    Args:
        file_path: Path to DOCX file or BytesIO object
        prefer_markitdown: Try MarkItDown first

    Returns:
        ExtractionResult with text, tables, and metadata
    """
    # Try MarkItDown first
    if prefer_markitdown:
        result = _try_markitdown_docx(file_path)
        if result:
            return result

    # Try python-docx
    return _try_python_docx(file_path)


def _try_markitdown_docx(file_path: Union[str, Path, BytesIO]) -> Optional[ExtractionResult]:
    """Try DOCX extraction with MarkItDown."""
    try:
        from markitdown import MarkItDown

        md = MarkItDown()

        # Handle BytesIO vs file path
        if isinstance(file_path, BytesIO):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                tmp.write(file_path.read())
                tmp_path = tmp.name
            file_path.seek(0)
            result = md.convert(tmp_path)
            Path(tmp_path).unlink()
        else:
            result = md.convert(str(file_path))

        if result and result.text_content:
            tables = _extract_markdown_tables(result.text_content)
            logger.info(f"MarkItDown extracted DOCX with {len(tables)} tables")
            return ExtractionResult(
                text=result.text_content,
                tables=tables,
                page_count=1,
                extraction_method="markitdown"
            )

    except ImportError:
        logger.debug("MarkItDown not available for DOCX")
    except Exception as e:
        logger.warning(f"MarkItDown DOCX extraction failed: {e}")

    return None


def _try_python_docx(file_path: Union[str, Path, BytesIO]) -> ExtractionResult:
    """Extract DOCX with python-docx."""
    try:
        from docx import Document

        doc = Document(file_path)
        content_parts = []
        tables = []

        # Extract paragraphs and tables in order
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)

        for i, table in enumerate(doc.tables):
            markdown = _docx_table_to_markdown(table)
            if markdown:
                content_parts.append(markdown)
                tables.append({"markdown": markdown, "index": i})

        logger.info(f"python-docx extracted {len(tables)} tables")
        return ExtractionResult(
            text="\n\n".join(content_parts),
            tables=tables,
            page_count=1,
            extraction_method="python-docx"
        )

    except ImportError:
        logger.error("python-docx not installed")
    except Exception as e:
        logger.error(f"python-docx extraction failed: {e}")

    return ExtractionResult(
        text="[DOCX extraction failed]",
        tables=[],
        page_count=0,
        extraction_method="failed"
    )


def _docx_table_to_markdown(table) -> Optional[str]:
    """Convert python-docx table to markdown."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace('|', '\\|')
            cells.append(cell_text)
        rows.append(cells)

    return _rows_to_markdown(rows)


# ============================================================================
# PPTX Extraction
# ============================================================================

def extract_pptx_with_tables(
    file_path: Union[str, Path, BytesIO],
    prefer_markitdown: bool = True
) -> ExtractionResult:
    """
    Extract text from PowerPoint with table awareness.

    Args:
        file_path: Path to PPTX file or BytesIO object
        prefer_markitdown: Try MarkItDown first

    Returns:
        ExtractionResult with text, tables, and metadata
    """
    # Try MarkItDown first
    if prefer_markitdown:
        result = _try_markitdown_pptx(file_path)
        if result:
            return result

    # Try python-pptx
    return _try_python_pptx(file_path)


def _try_markitdown_pptx(file_path: Union[str, Path, BytesIO]) -> Optional[ExtractionResult]:
    """Try PPTX extraction with MarkItDown."""
    try:
        from markitdown import MarkItDown

        md = MarkItDown()

        if isinstance(file_path, BytesIO):
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(file_path.read())
                tmp_path = tmp.name
            file_path.seek(0)
            result = md.convert(tmp_path)
            Path(tmp_path).unlink()
        else:
            result = md.convert(str(file_path))

        if result and result.text_content:
            tables = _extract_markdown_tables(result.text_content)
            logger.info(f"MarkItDown extracted PPTX with {len(tables)} tables")
            return ExtractionResult(
                text=result.text_content,
                tables=tables,
                page_count=1,
                extraction_method="markitdown"
            )

    except ImportError:
        logger.debug("MarkItDown not available for PPTX")
    except Exception as e:
        logger.warning(f"MarkItDown PPTX extraction failed: {e}")

    return None


def _try_python_pptx(file_path: Union[str, Path, BytesIO]) -> ExtractionResult:
    """Extract PPTX with python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        content_parts = []
        slide_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_count += 1
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            if slide_text:
                content_parts.append(
                    f"[Slide {slide_num}]\n" + "\n".join(slide_text)
                )

        logger.info(f"python-pptx extracted {slide_count} slides")
        return ExtractionResult(
            text="\n\n".join(content_parts),
            tables=[],
            page_count=slide_count,
            extraction_method="python-pptx"
        )

    except ImportError:
        logger.error("python-pptx not installed")
    except Exception as e:
        logger.error(f"python-pptx extraction failed: {e}")

    return ExtractionResult(
        text="[PPTX extraction failed]",
        tables=[],
        page_count=0,
        extraction_method="failed"
    )


# ============================================================================
# Helper Functions
# ============================================================================

def _table_to_markdown(table: List[List[str]]) -> str:
    """Convert 2D array to markdown table."""
    if not table or len(table) < 2:
        return ""

    # Clean and normalize
    rows = []
    for row in table:
        cells = [str(cell).strip().replace('|', '\\|') if cell else "" for cell in row]
        rows.append(cells)

    return _rows_to_markdown(rows)


def _rows_to_markdown(rows: List[List[str]]) -> Optional[str]:
    """Convert rows to markdown table format."""
    if len(rows) < 2:
        return None

    # Normalize column count
    max_cols = max(len(row) for row in rows)
    for row in rows:
        while len(row) < max_cols:
            row.append("")

    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _extract_markdown_tables(text: str) -> List[Dict[str, Any]]:
    """Extract markdown tables from text."""
    import re

    tables = []
    # Pattern to match markdown tables
    table_pattern = r'\|[^\n]+\|\n\|[-:\s|]+\|\n(?:\|[^\n]+\|\n?)+'

    for i, match in enumerate(re.finditer(table_pattern, text)):
        tables.append({
            "markdown": match.group(0),
            "index": i,
            "start": match.start(),
            "end": match.end()
        })

    return tables


# ============================================================================
# Unified Extraction
# ============================================================================

def extract_document(
    file_path: Union[str, Path],
    file_type: Optional[str] = None,
    prefer_markitdown: bool = True
) -> ExtractionResult:
    """
    Extract text from any supported document type.

    Args:
        file_path: Path to the document
        file_type: File type override (pdf, docx, pptx, txt)
        prefer_markitdown: Try MarkItDown first

    Returns:
        ExtractionResult with text, tables, and metadata
    """
    path = Path(file_path)
    ext = file_type or path.suffix.lower().lstrip('.')

    if ext == 'pdf':
        return extract_pdf_with_tables(file_path, prefer_markitdown)
    elif ext in ('docx', 'doc'):
        return extract_docx_with_tables(file_path, prefer_markitdown)
    elif ext in ('pptx', 'ppt'):
        return extract_pptx_with_tables(file_path, prefer_markitdown)
    elif ext in ('txt', 'md'):
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return ExtractionResult(
            text=text,
            tables=_extract_markdown_tables(text),
            page_count=1,
            extraction_method="plaintext"
        )
    else:
        logger.warning(f"Unsupported file type: {ext}")
        return ExtractionResult(
            text=f"[Unsupported file type: {ext}]",
            tables=[],
            page_count=0,
            extraction_method="unsupported"
        )
